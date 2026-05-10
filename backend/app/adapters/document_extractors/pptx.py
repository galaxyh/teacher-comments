"""PptxExtractor — python-pptx, slide-per-section.

Per DESIGN-001 §8.2: each slide → markdown section with title (if any) and
text frames. Embedded images set `has_images=True` so ProcessingPipeline can
optionally route to a vision pass (Phase 7+).
"""

from __future__ import annotations

import asyncio
import io
import logging
from typing import Final

from pptx import Presentation
from pptx.util import Emu

from app.adapters.document_extractors import ExtractionResult
from app.core.exceptions import DocumentExtractionError, UnsupportedFormatError

logger = logging.getLogger(__name__)

PPTX_MIME: Final[str] = (
    "application/vnd.openxmlformats-officedocument.presentationml.presentation"
)


class PptxExtractor:
    def supports(self, *, mime_type: str, filename: str) -> bool:
        return mime_type == PPTX_MIME or filename.lower().endswith(".pptx")

    async def extract(self, *, file_bytes: bytes, filename: str) -> ExtractionResult:
        if file_bytes[:4] == b"\xd0\xcf\x11\xe0":
            raise UnsupportedFormatError(
                "Legacy .ppt (OLE) — convert to .pptx",
                context={"filename": filename},
            )
        try:
            return await asyncio.to_thread(self._sync_extract, file_bytes, filename)
        except UnsupportedFormatError:
            raise
        except Exception as exc:  # noqa: BLE001
            raise DocumentExtractionError(
                f"Failed to parse .pptx: {exc.__class__.__name__}: {exc}",
                context={"filename": filename},
            ) from exc

    def _sync_extract(self, file_bytes: bytes, filename: str) -> ExtractionResult:
        prs = Presentation(io.BytesIO(file_bytes))
        sections: list[str] = []
        has_images = False
        warnings: list[str] = []

        for i, slide in enumerate(prs.slides, start=1):
            title = ""
            body_lines: list[str] = []
            for shape in slide.shapes:
                # python-pptx exposes 'has_text_frame' and 'shape_type' as duck-typed attrs
                if getattr(shape, "shape_type", None) == 13:  # 13 = picture
                    has_images = True
                if not shape.has_text_frame:
                    continue
                for para in shape.text_frame.paragraphs:
                    text = "".join(run.text for run in para.runs).strip()
                    if not text:
                        continue
                    # First non-empty line of the title placeholder = slide title
                    if (
                        not title
                        and shape.has_text_frame
                        and shape.is_placeholder
                        and shape.placeholder_format
                        and shape.placeholder_format.idx == 0
                    ):
                        title = text
                    else:
                        body_lines.append(text)

            heading = f"## 第 {i} 張投影片" + (f"：{title}" if title else "")
            body = "\n".join(f"- {line}" for line in body_lines) if body_lines else ""
            sections.append(f"{heading}\n{body}".strip())

        if not sections:
            warnings.append("presentation had no slides with text")
        return ExtractionResult(
            text="\n\n".join(sections).strip(),
            has_images=has_images,
            page_count=len(prs.slides),
            warnings=warnings,
        )


# Suppress unused-import warning — `Emu` import surfaces "missing pptx" earlier
_ = Emu
