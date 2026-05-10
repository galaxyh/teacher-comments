"""Tests for the Phase 5b extractors: xlsx, pptx, pdf."""

from __future__ import annotations

import io

import pytest
from openpyxl import Workbook
from pptx import Presentation
from pptx.util import Inches
from pypdf import PdfWriter

from app.adapters.document_extractors.pdf import PdfExtractor
from app.adapters.document_extractors.pptx import PptxExtractor
from app.adapters.document_extractors.xlsx import XlsxExtractor
from app.core.exceptions import DocumentExtractionError, UnsupportedFormatError


# ── XLSX ─────────────────────────────────────────────────────────


def _build_xlsx(*, rows: list[list[str]], sheet_name: str = "Sheet1") -> bytes:
    wb = Workbook()
    ws = wb.active
    ws.title = sheet_name
    for row in rows:
        ws.append(row)
    buf = io.BytesIO()
    wb.save(buf)
    return buf.getvalue()


class TestXlsxExtractor:
    @pytest.mark.asyncio
    async def test_basic_table(self) -> None:
        ext = XlsxExtractor()
        body = _build_xlsx(
            rows=[["姓名", "成績"], ["小明", "90"], ["小華", "85"]],
            sheet_name="期中",
        )
        result = await ext.extract(file_bytes=body, filename="grades.xlsx")
        assert "## 期中" in result.text
        assert "| 姓名 | 成績 |" in result.text
        assert "| 小明 | 90 |" in result.text
        assert "| 小華 | 85 |" in result.text

    @pytest.mark.asyncio
    async def test_blank_sheet_emits_warning(self) -> None:
        ext = XlsxExtractor()
        body = _build_xlsx(rows=[], sheet_name="empty")
        result = await ext.extract(file_bytes=body, filename="x.xlsx")
        assert any("empty" in w for w in result.warnings)

    @pytest.mark.asyncio
    async def test_legacy_xls_unsupported(self) -> None:
        ext = XlsxExtractor()
        legacy = b"\xd0\xcf\x11\xe0" + b"\x00" * 100
        with pytest.raises(UnsupportedFormatError, match="Legacy"):
            await ext.extract(file_bytes=legacy, filename="old.xls")

    @pytest.mark.asyncio
    async def test_corrupt_xlsx_raises_extraction_error(self) -> None:
        ext = XlsxExtractor()
        with pytest.raises(DocumentExtractionError):
            await ext.extract(file_bytes=b"not a real xlsx", filename="x.xlsx")


# ── PPTX ─────────────────────────────────────────────────────────


def _build_pptx(slides: list[tuple[str | None, list[str]]]) -> bytes:
    prs = Presentation()
    blank_layout = prs.slide_layouts[5]  # title + content
    for title_text, body_texts in slides:
        slide = prs.slides.add_slide(blank_layout)
        if title_text and slide.shapes.title is not None:
            slide.shapes.title.text = title_text
        # Add a body text box
        if body_texts:
            box = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(6), Inches(3))
            tf = box.text_frame
            for i, text in enumerate(body_texts):
                p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
                p.text = text
    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


class TestPptxExtractor:
    @pytest.mark.asyncio
    async def test_basic_slide(self) -> None:
        ext = PptxExtractor()
        body = _build_pptx([("第一張", ["重點 A", "重點 B"])])
        result = await ext.extract(file_bytes=body, filename="t.pptx")
        assert "第 1 張投影片" in result.text
        assert "第一張" in result.text  # title appended
        assert "重點 A" in result.text and "重點 B" in result.text

    @pytest.mark.asyncio
    async def test_multiple_slides_become_sections(self) -> None:
        ext = PptxExtractor()
        body = _build_pptx([
            ("Intro", ["bullet 1"]),
            ("Body", ["bullet 2"]),
        ])
        result = await ext.extract(file_bytes=body, filename="t.pptx")
        assert "第 1 張投影片" in result.text and "第 2 張投影片" in result.text
        assert result.page_count == 2

    @pytest.mark.asyncio
    async def test_legacy_ppt_unsupported(self) -> None:
        ext = PptxExtractor()
        with pytest.raises(UnsupportedFormatError, match="Legacy"):
            await ext.extract(file_bytes=b"\xd0\xcf\x11\xe0\x00", filename="old.ppt")


# ── PDF ──────────────────────────────────────────────────────────


def _build_pdf(*, encrypted: bool = False) -> bytes:
    """Build a minimal PDF using pypdf. Empty PDFs lack text by design — use
    fpdf2 for richer fixtures if needed (out of scope here; we test plumbing)."""
    writer = PdfWriter()
    writer.add_blank_page(width=612, height=792)
    if encrypted:
        writer.encrypt(user_password="locked", owner_password="locked")
    buf = io.BytesIO()
    writer.write(buf)
    return buf.getvalue()


class TestPdfExtractor:
    @pytest.mark.asyncio
    async def test_blank_pdf_no_text_returns_warning(self) -> None:
        ext = PdfExtractor()
        result = await ext.extract(file_bytes=_build_pdf(), filename="blank.pdf")
        # Blank page → empty extracted text; warning surfaces 'no extractable text'
        assert result.text == ""
        assert any("no extractable text" in w for w in result.warnings)
        assert result.page_count == 1

    @pytest.mark.asyncio
    async def test_encrypted_pdf_unsupported(self) -> None:
        ext = PdfExtractor()
        body = _build_pdf(encrypted=True)
        with pytest.raises(UnsupportedFormatError, match="Encrypted"):
            await ext.extract(file_bytes=body, filename="locked.pdf")

    @pytest.mark.asyncio
    async def test_corrupt_pdf_raises_extraction_error(self) -> None:
        ext = PdfExtractor()
        with pytest.raises(DocumentExtractionError):
            await ext.extract(file_bytes=b"not a real pdf", filename="bad.pdf")


# ── Registry routing ─────────────────────────────────────────────


def test_registry_routes_to_correct_extractor() -> None:
    from app.adapters.document_extractors import build_default_registry
    from app.adapters.document_extractors.docx import DocxExtractor

    reg = build_default_registry()
    # XLSX MIME
    assert isinstance(
        reg.get(
            mime_type="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
            filename="x",
        ),
        XlsxExtractor,
    )
    # PPTX
    assert isinstance(
        reg.get(
            mime_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
            filename="x",
        ),
        PptxExtractor,
    )
    # PDF
    assert isinstance(reg.get(mime_type="application/pdf", filename="x"), PdfExtractor)
    # DOCX still routes correctly (no regression)
    assert isinstance(
        reg.get(
            mime_type="application/vnd.openxmlformats-officedocument.wordprocessingml.document",
            filename="x",
        ),
        DocxExtractor,
    )
